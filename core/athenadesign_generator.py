import os
import re
import json
import httpx

# System instructions to guide the LLM to think like Claude Design / Open Design
SYSTEM_PROMPT = """You are AthenaDesign, an elite AI design companion. Your goal is to generate beautiful, interactive, premium visual interfaces or Python scripts based on user requests.

⚠️ ABSOLUTE PRIORITY — THE CODE MUST RUN (this overrides every styling instruction below):
- Output SYNTACTICALLY VALID code. CLOSE EVERY tag correctly (`</span>`, `</div>` — NEVER `</}` or a stray `}`), balance every brace/parenthesis/bracket, and close every string.
- NEVER truncate. Output the COMPLETE code from first to last line. If it would be too long, build something smaller but COMPLETE.
- For ```jsx: EXACTLY ONE component named `App`. NO `import`, NO `export` (React, ReactDOM, hooks and Tailwind are already global). Do NOT call ReactDOM yourself.
- A beautiful design that does not compile is a FAILURE. Correctness first, then beauty.

OUTPUT FORMAT (follow EXACTLY):
1. First, a SHORT explanation in French (2 to 4 sentences) of your design choices.
2. Then the COMPLETE code inside ONE single fenced block tagged with the language:
   ```html      for web designs (static HTML/CSS/JS)
   ```python    for scripts / PowerPoint / charts
   ```jsx       for React components / interactive apps
   ```mermaid   for diagrams (flowchart, sequence, class, ER, gantt, state, mindmap…)
   Put ALL the code in that single block. Never truncate. Do NOT mix explanation text inside the code. Do NOT output two code blocks.
   ── MULTI-FILE PROJECTS (preferred for anything beyond a tiny standalone page) ──
   For a STRUCTURED project (real site with separate CSS/JS, multiple pages, assets, or a React
   app split into files), output SEVERAL files instead of one block. Precede EACH file with a line:
       === FILE: relative/path ===
   then its content (optionally in a fenced block). Rules:
   - The entry point MUST be `index.html`. Keep the French explanation BEFORE the first `=== FILE: ===`.
   - Use RELATIVE links between files: `<link rel="stylesheet" href="./css/style.css">`,
     `<script src="./js/app.js"></script>`, images in `./assets/…`. NEVER absolute paths.
   - Typical layout: `index.html`, `css/style.css`, `js/app.js` (+ extra pages/components as needed).
   - Keep a SINGLE fenced block only for a genuinely simple, one-file artifact, a diagram, or Python.
3. Finally, under the code block, you MUST append:
   - A `<suggestions>` block containing 3 proposed next steps (one per line, prefixed by a dash `-`).
   - For web layouts (html/jsx), a `<tweaks>` block containing 2 to 4 visual controls mapping to CSS variables used in your layout.
     Format of each tweak line: `type | Label | variableName | rangeOrValues | defaultValue`
     Where:
       - `type` is one of: `color` (color picker), `slider` (range), `toggle` (true/false switch), or `select` (dropdown)
       - `rangeOrValues` is `min,max` for sliders (ex: `10px,50px`), `val1,val2,val3` for select (ex: `flat,neumorph,glass`), or empty for color/toggle
     Examples:
       - `color | Boutons CTA | --btn-color | | #e11d48`
       - `slider | Espacement | --section-gap | 20px,120px | 60px`
       - `toggle | Mode Sombre | --dark-theme | | false`
       - `select | Style Cartes | --card-style | flat,glass,neobrutal | glass`
     ⚠️ CRITICAL — a tweak is USELESS unless the CSS reads its variable. For EVERY tweak you declare you MUST:
       1) define the variable on `:root` with its default value, e.g. `:root{ --btn-color:#e11d48; --section-gap:60px; }`
       2) ACTUALLY USE it everywhere relevant via `var(--name)` — NEVER hard-code that value elsewhere.
          ex: `.cta{ background: var(--btn-color); }`, `section{ padding: var(--section-gap); }`.
       For React/Tailwind: still drive the tweakable properties through these CSS variables (inline style
       `style={{background:'var(--btn-color)'}}` or a `<style>:root{…}</style>` block + `var(--…)`), NOT via
       hard-coded Tailwind color classes — otherwise the right-side controls change nothing.
Choose REACT (```jsx) for stateful / interactive UIs (forms with logic, tabs, todo, calculators, dashboards with interactivity). Choose ```html for static/visual pages. Choose ```python for data, charts or .pptx.

REACT RULES (when you output ```jsx):
- Define a single component named `App`. Do NOT write `import` or `export` statements — React,
  ReactDOM, the hooks and Tailwind CSS are ALREADY provided globally in the page.
- Use hooks via `React.useState`, `React.useEffect`, etc. (or destructure: `const {useState}=React;`).
- Style with Tailwind utility classes (available) and/or inline styles. Make it premium and responsive.
- The page renders <App/> automatically; do not call ReactDOM yourself.

MERMAID RULES (when you output ```mermaid): write PURE Mermaid syntax only (no HTML, no
markdown), starting with the diagram type, e.g. `flowchart TD`, `sequenceDiagram`,
`classDiagram`, `erDiagram`, `stateDiagram-v2`, `gantt`, `mindmap`. The page renders it.

DESIGN RULES:
1. VISUAL DIVERSITY & THEME ADAPTATION: Do NOT always produce the same dark-mode glassmorphic style. Be highly creative and tailor the visual design entirely to the topic, industry, and tone of the request:
   - Corporate/Professional: clean, light-mode layouts, generous whitespace, structured grid systems, trustworthy colors (indigo, slate, sky, emerald).
   - Creative/Portfolio: bold typography, unique asymmetric layouts, high-contrast, neobrutalism, or artistic flat styles.
   - Minimalist/Editorial: gorgeous large serif or clean sans-serif headings, subtle borders, black-and-white accents with one primary brand color (e.g. amber, teal).
   - Tech/Futuristic: modern dark mode, neon glow accents, cyber-gradients.
   Select matching Google Fonts (e.g., Playfair Display, Montserrat, Outfit, Inter, Space Grotesk, Syne) to match the selected aesthetic.
2. VISUAL EXCELLENCE: Never use plain basic colors (red, blue, green). Use curated palettes with harmonious hues.
3. INTERACTIVITY & MICRO-ANIMATIONS: Implement responsive state changes (hover scales, shadow shifts, active press states, smooth list transition keyframes) so the interface feels alive and reactive.
4. SELF-CONTAINED HTML: a web artifact is ONE standalone .html file. If you use a library, you MUST load it from a CDN in the file. In particular, if you use Lucide icons (<i data-lucide="...">), you MUST add `<script src="https://unpkg.com/lucide@latest"></script>` AND call `lucide.createIcons()` after the DOM is ready — otherwise icons stay invisible. Same for Chart.js / FontAwesome. If unsure, prefer inline SVG or emoji. Never reference a library you did not load.
5. PYTHON CODE: generate useful data, charts (Matplotlib/Plotly, clean modern styling, no grey background) or PowerPoint via python-pptx saved to the current directory. End charts with plt.show()/fig.show() to render a preview.
6. POWERPOINT — SIMPLE, ROBUST & READABLE (critical). python-pptx is NOT HTML/CSS: fancy custom
   layouts turn into an unreadable mess. Follow these rules strictly:
   - PREFER BUILT-IN LAYOUTS & PLACEHOLDERS: `prs.slide_layouts[0]` (title) and `[1]` (title+content).
     Put the title in `slide.shapes.title` and bullets in the content placeholder. AVOID piling up
     many manually-positioned text boxes and decorative rectangles.
   - NEVER use negative or zero width/height. NEVER draw charts/bars by hand with shapes (no “bars”
     via rectangles, no negative heights). If you truly need a chart, use a Matplotlib PNG inserted
     via `slide.shapes.add_picture`, OR just describe the data as bullets. Keep it simple.
   - STAY IN BOUNDS: slide is 13.333 in × 7.5 in. Every shape: left≥0, top≥0, left+width ≤ 13.0,
     top+height ≤ 7.0. ≤ 5-6 short bullets/slide; split long content across slides.
   - READABLE COLORS WITH CONTRAST: set the slide BACKGROUND on every slide
     (`slide.background.fill.solid(); slide.background.fill.fore_color.rgb = COLOR_BG`) AND set an
     EXPLICIT, CONTRASTING font color on the title and on EVERY body run
     (`run.font.color.rgb = COLOR_TEXT`). Light text on dark bg OR dark text on light bg — NEVER rely
     on the default ('auto') color and NEVER dark-on-dark. Define a small RGBColor palette once and reuse it.
   - LAYERING: add any decorative/colored shape (rectangle, band, accent) BEFORE the text it
     decorates, so the TEXT stays ON TOP. A colored shape must NEVER cover text. Prefer placing
     decorations in empty zones (header band, side strip) rather than behind paragraphs.
   - Body frames: `tf.word_wrap = True`; titles 32-40pt, body 16-20pt.
7. RESPONSIVENESS (HTML): fluid layouts (%, vw/vh, flexbox, CSS grid auto-fit, fr units), not fixed pixel widths. Must scale to tablet/phone without horizontal scroll.
"""

# Premium mock templates to run offline
MOCK_TEMPLATES = {
    "dashboard": {
        "type": "html",
        "explanation": "Voici un tableau de bord analytique moderne avec des graphiques interactifs (via Chart.js), un thème sombre type Glassmorphism, et des effets de survol fluides.",
        "code": """<!DOCTYPE html>
<html lang="fr">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Sleek Analytics Dashboard</title>
    <link href="https://fonts.googleapis.com/css2?family=Outfit:wght@300;400;600;800&display=swap" rel="stylesheet">
    <script src="https://cdn.jsdelivr.net/npm/chart.js"></script>
    <style>
        :root {
            --bg-color: #0b0f19;
            --panel-bg: rgba(20, 26, 45, 0.6);
            --border-color: rgba(255, 255, 255, 0.08);
            --text-primary: #f8fafc;
            --text-secondary: #94a3b8;
            --accent-primary: #6366f1;
            --accent-secondary: #d946ef;
            --accent-success: #10b981;
        }

        * {
            box-sizing: border-box;
            margin: 0;
            padding: 0;
            font-family: 'Outfit', sans-serif;
            transition: all 0.3s ease;
        }

        body {
            background: radial-gradient(circle at top right, #1e1b4b, var(--bg-color));
            color: var(--text-primary);
            min-height: 100vh;
            display: flex;
            padding: 24px;
            overflow-x: hidden;
        }

        .container {
            width: 100%;
            max-width: 1400px;
            margin: 0 auto;
            display: grid;
            grid-template-columns: 240px 1fr;
            gap: 24px;
        }

        /* Sidebar Glassmorphism */
        sidebar {
            background: var(--panel-bg);
            backdrop-filter: blur(12px);
            border: 1px solid var(--border-color);
            border-radius: 20px;
            padding: 24px;
            display: flex;
            flex-direction: column;
            gap: 32px;
        }

        .logo {
            font-size: 1.5rem;
            font-weight: 800;
            background: linear-gradient(135deg, var(--accent-primary), var(--accent-secondary));
            -webkit-background-clip: text;
            -webkit-text-fill-color: transparent;
            letter-spacing: 1px;
        }

        .menu {
            list-style: none;
            display: flex;
            flex-direction: column;
            gap: 12px;
        }

        .menu-item {
            padding: 12px 16px;
            border-radius: 12px;
            color: var(--text-secondary);
            cursor: pointer;
            font-weight: 500;
        }

        .menu-item:hover, .menu-item.active {
            background: rgba(99, 102, 241, 0.15);
            color: var(--text-primary);
            border-left: 4px solid var(--accent-primary);
            padding-left: 20px;
        }

        /* Main Content */
        main {
            display: flex;
            flex-direction: column;
            gap: 24px;
        }

        header {
            display: flex;
            justify-content: space-between;
            align-items: center;
        }

        h1 {
            font-size: 2rem;
            font-weight: 600;
        }

        .kpi-grid {
            display: grid;
            grid-template-columns: repeat(auto-fit, minmax(240px, 1fr));
            gap: 20px;
        }

        .kpi-card {
            background: var(--panel-bg);
            backdrop-filter: blur(12px);
            border: 1px solid var(--border-color);
            border-radius: 16px;
            padding: 24px;
            position: relative;
            overflow: hidden;
        }

        .kpi-card::before {
            content: '';
            position: absolute;
            top: 0;
            left: 0;
            width: 100%;
            height: 4px;
            background: linear-gradient(90deg, var(--accent-primary), var(--accent-secondary));
            opacity: 0;
        }

        .kpi-card:hover {
            transform: translateY(-5px);
            box-shadow: 0 10px 20px rgba(0,0,0,0.3);
            border-color: rgba(99, 102, 241, 0.3);
        }

        .kpi-card:hover::before {
            opacity: 1;
        }

        .kpi-label {
            color: var(--text-secondary);
            font-size: 0.9rem;
            margin-bottom: 8px;
        }

        .kpi-value {
            font-size: 1.8rem;
            font-weight: 700;
            margin-bottom: 8px;
        }

        .kpi-change {
            font-size: 0.85rem;
            color: var(--accent-success);
            display: flex;
            align-items: center;
            gap: 4px;
        }

        .kpi-change.down {
            color: #ef4444;
        }

        /* Charts Section */
        .chart-section {
            display: grid;
            grid-template-columns: 2fr 1fr;
            gap: 20px;
        }

        .chart-card {
            background: var(--panel-bg);
            backdrop-filter: blur(12px);
            border: 1px solid var(--border-color);
            border-radius: 20px;
            padding: 24px;
        }

        .chart-header {
            margin-bottom: 20px;
            display: flex;
            justify-content: space-between;
            align-items: center;
        }

        .chart-title {
            font-size: 1.1rem;
            font-weight: 600;
        }

        /* Responsive Media Queries */
        @media (max-width: 900px) {
            .container {
                grid-template-columns: 1fr;
            }
            sidebar {
                flex-direction: row;
                justify-content: space-between;
                align-items: center;
                padding: 16px;
            }
            .menu {
                flex-direction: row;
                gap: 12px;
            }
            .chart-section {
                grid-template-columns: 1fr;
            }
        }
        @media (max-width: 600px) {
            sidebar {
                flex-direction: column;
                gap: 16px;
            }
            .menu {
                flex-wrap: wrap;
                justify-content: center;
            }
            body {
                padding: 12px;
            }
        }
    </style>
</head>
<body>
    <div class="container">
        <sidebar>
            <div class="logo">PyDesign OS</div>
            <ul class="menu">
                <li class="menu-item active">Dashboard</li>
                <li class="menu-item">Analyses</li>
                <li class="menu-item">Projets</li>
                <li class="menu-item">Configuration</li>
            </ul>
        </sidebar>
        
        <main>
            <header>
                <div>
                    <h1>Aperçu Analytique</h1>
                    <p style="color: var(--text-secondary); margin-top: 4px;">Données en temps réel de votre application</p>
                </div>
                <div class="user-pill">Admin Mode</div>
            </header>

            <div class="kpi-grid">
                <div class="kpi-card">
                    <div class="kpi-label">Utilisateurs Actifs</div>
                    <div class="kpi-value">12,842</div>
                    <div class="kpi-change">▲ +12.4% cette semaine</div>
                </div>
                <div class="kpi-card">
                    <div class="kpi-label">Revenus Mensuels</div>
                    <div class="kpi-value">€45,210</div>
                    <div class="kpi-change">▲ +8.1% vs mois dernier</div>
                </div>
                <div class="kpi-card">
                    <div class="kpi-label">Temps de Réponse</div>
                    <div class="kpi-value">84ms</div>
                    <div class="kpi-change down">▼ -2.4% (amélioration)</div>
                </div>
            </div>

            <div class="chart-section">
                <div class="chart-card">
                    <div class="chart-header">
                        <span class="chart-title">Croissance des utilisateurs</span>
                    </div>
                    <canvas id="lineChart" height="180"></canvas>
                </div>
                <div class="chart-card">
                    <div class="chart-header">
                        <span class="chart-title">Canaux d'acquisition</span>
                    </div>
                    <canvas id="doughnutChart" height="180"></canvas>
                </div>
            </div>
        </main>
    </div>

    <script>
        // Line Chart Initialization
        const ctxLine = document.getElementById('lineChart').getContext('2d');
        new Chart(ctxLine, {
            type: 'line',
            data: {
                labels: ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun'],
                datasets: [{
                    label: 'Utilisateurs',
                    data: [4000, 5500, 7000, 6500, 9000, 12842],
                    borderColor: '#6366f1',
                    backgroundColor: 'rgba(99, 102, 241, 0.1)',
                    borderWidth: 3,
                    fill: true,
                    tension: 0.4
                }]
            },
            options: {
                responsive: true,
                plugins: { legend: { display: false } },
                scales: {
                    x: { grid: { display: false }, ticks: { color: '#94a3b8' } },
                    y: { grid: { color: 'rgba(255,255,255,0.05)' }, ticks: { color: '#94a3b8' } }
                }
            }
        });

        // Doughnut Chart Initialization
        const ctxDoughnut = document.getElementById('doughnutChart').getContext('2d');
        new Chart(ctxDoughnut, {
            type: 'doughnut',
            data: {
                labels: ['Recherche', 'Direct', 'Réseaux', 'Parrainage'],
                datasets: [{
                    data: [45, 25, 20, 10],
                    backgroundColor: ['#6366f1', '#d946ef', '#10b981', '#f59e0b'],
                    borderWidth: 0
                }]
            },
            options: {
                responsive: true,
                plugins: {
                    legend: {
                        position: 'bottom',
                        labels: { color: '#94a3b8', font: { family: 'Outfit' } }
                    }
                }
            }
        });
    </script>
</body>
</html>"""
    },
    "plot": {
        "type": "python",
        "explanation": "Voici un script Python qui génère des courbes de sinus et de cosinus bruitées avec un style sombre moderne (grille indigo/slate) et calcule des moyennes mobiles via Pandas. Cliquez sur l'onglet **Preview** pour voir le graphique généré !",
        "code": """import matplotlib.pyplot as plt
import numpy as np
import pandas as pd

# Génération de données temporelles
np.random.seed(42)
x = np.linspace(0, 10, 100)
y_sin = np.sin(x) + np.random.normal(0, 0.1, 100)
y_cos = np.cos(x) + np.random.normal(0, 0.1, 100)

df = pd.DataFrame({
    'x': x,
    'Sinus': y_sin,
    'Cosinus': y_cos
})

# Calcul d'une moyenne mobile
df['Sinus_smooth'] = df['Sinus'].rolling(window=5, min_periods=1).mean()

# Configuration du style sombre pour le graphique
plt.style.use('dark_background')
fig, ax = plt.subplots(figsize=(10, 5), dpi=150)

# Tracé des courbes avec des couleurs harmonieuses
ax.plot(df['x'], df['Sinus'], label='Sinus bruité', color='#6366f1', alpha=0.4, linestyle='--')
ax.plot(df['x'], df['Sinus_smooth'], label='Moyenne mobile (5)', color='#6366f1', linewidth=2.5)
ax.plot(df['x'], df['Cosinus'], label='Cosinus bruité', color='#d946ef', alpha=0.7, linewidth=1.5)

# Styling des axes et légendes
ax.spines['top'].set_visible(False)
ax.spines['right'].set_visible(False)
ax.spines['left'].set_color('#334155')
ax.spines['bottom'].set_color('#334155')

ax.xaxis.grid(True, linestyle=':', alpha=0.3, color='#475569')
ax.yaxis.grid(True, linestyle=':', alpha=0.3, color='#475569')

ax.set_title("Analyse de signaux avec moyennes mobiles", fontsize=14, pad=15, weight='bold', color='#f8fafc')
ax.set_xlabel("Temps (s)", color='#94a3b8')
ax.set_ylabel("Amplitude", color='#94a3b8')

# Custom legend
ax.legend(facecolor='#0f172a', edgecolor='#1e293b', labelcolor='#e2e8f0')

print("Analyse de signal terminée.")
print(f"Moyenne du sinus : {df['Sinus'].mean():.4f}")
print(f"Moyenne du cosinus : {df['Cosinus'].mean():.4f}")

plt.show()
"""
    },
    "simulation": {
        "type": "python",
        "explanation": "Voici une simulation physique d'un attracteur chaotique de Lorenz (théorie du chaos) tracée en 3D avec Matplotlib. Le script résout les équations différentielles et les affiche sur un graphique tridimensionnel.",
        "code": """import numpy as np
import matplotlib.pyplot as plt

def lorenz(x, y, z, s=10, r=28, b=2.667):
    x_dot = s*(y - x)
    y_dot = r*x - y - x*z
    z_dot = x*y - b*z
    return x_dot, y_dot, z_dot

dt = 0.01
num_steps = 2500

xs = np.empty(num_steps + 1)
ys = np.empty(num_steps + 1)
zs = np.empty(num_steps + 1)

xs[0], ys[0], zs[0] = (0., 1., 1.05)

for i in range(num_steps):
    x_dot, y_dot, z_dot = lorenz(xs[i], ys[i], zs[i])
    xs[i + 1] = xs[i] + (x_dot * dt)
    ys[i + 1] = ys[i] + (y_dot * dt)
    zs[i + 1] = zs[i] + (z_dot * dt)

# Tracé 3D
plt.style.use('dark_background')
fig = plt.figure(figsize=(10, 7), dpi=150)
ax = fig.add_subplot(projection='3d')

# Couleur dégradée le long de la trajectoire
colors = plt.cm.plasma(np.linspace(0, 1, num_steps))

for i in range(num_steps):
    ax.plot(xs[i:i+2], ys[i:i+2], zs[i:i+2], color=colors[i], alpha=0.8, linewidth=0.7)

ax.set_title("Attracteur chaotique de Lorenz (Simulation 3D)", fontsize=14, pad=15, weight='bold', color='#f8fafc')
ax.axis('off') # masquer les axes pour un rendu artistique pur

print("Attracteur de Lorenz simulé avec succès.")
print(f"Coordonnées finales : ({xs[-1]:.2f}, {ys[-1]:.2f}, {zs[-1]:.2f})")

plt.show()
"""
    },
    "form": {
        "type": "html",
        "explanation": "Voici un formulaire de contact/feedback moderne avec un effet de flou d'arrière-plan (glassmorphism), des micro-animations sur les boutons de validation, et des dégradés subtils.",
        "code": """<!DOCTYPE html>
<html lang="fr">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Feedback Hub</title>
    <link href="https://fonts.googleapis.com/css2?family=Outfit:wght@300;400;600;700&display=swap" rel="stylesheet">
    <style>
        body {
            background: linear-gradient(135deg, #0f172a 0%, #1e1b4b 100%);
            color: #f8fafc;
            min-height: 100vh;
            display: flex;
            align-items: center;
            justify-content: center;
            font-family: 'Outfit', sans-serif;
            margin: 0;
            padding: 20px;
        }
        .form-card {
            background: rgba(30, 41, 59, 0.4);
            backdrop-filter: blur(16px);
            border: 1px solid rgba(255, 255, 255, 0.08);
            border-radius: 24px;
            padding: 40px;
            width: 100%;
            max-width: 480px;
            box-shadow: 0 20px 40px rgba(0,0,0,0.3);
            text-align: center;
            transition: transform 0.3s ease;
        }
        .form-card:hover {
            transform: translateY(-2px);
        }
        h2 {
            font-size: 2rem;
            margin-bottom: 8px;
            background: linear-gradient(135deg, #6366f1, #d946ef);
            -webkit-background-clip: text;
            -webkit-text-fill-color: transparent;
        }
        p {
            color: #94a3b8;
            margin-bottom: 32px;
            font-size: 0.95rem;
        }
        .input-group {
            text-align: left;
            margin-bottom: 20px;
        }
        label {
            display: block;
            font-size: 0.85rem;
            font-weight: 600;
            color: #cbd5e1;
            margin-bottom: 8px;
            letter-spacing: 0.5px;
        }
        input, textarea {
            width: 100%;
            background: rgba(15, 23, 42, 0.6);
            border: 1px solid rgba(255,255,255,0.1);
            border-radius: 12px;
            padding: 12px 16px;
            color: white;
            font-size: 0.95rem;
            box-sizing: border-box;
            transition: all 0.2s;
        }
        input:focus, textarea:focus {
            outline: none;
            border-color: #6366f1;
            box-shadow: 0 0 0 3px rgba(99, 102, 241, 0.2);
        }
        button {
            width: 100%;
            background: linear-gradient(135deg, #6366f1 0%, #d946ef 100%);
            color: white;
            border: none;
            padding: 14px;
            border-radius: 12px;
            font-size: 1rem;
            font-weight: 600;
            cursor: pointer;
            box-shadow: 0 4px 12px rgba(99, 102, 241, 0.3);
            margin-top: 10px;
            transition: all 0.2s;
        }
        button:hover {
            transform: translateY(-1px);
            box-shadow: 0 6px 20px rgba(99, 102, 241, 0.4);
        }
        
        /* Responsive Media Query */
        @media (max-width: 480px) {
            .form-card {
                padding: 24px 16px;
                border-radius: 16px;
            }
            body {
                padding: 12px;
            }
        }
    </style>
</head>
<body>
    <div class="form-card">
        <h2>Nous Contacter</h2>
        <p>Envoyez-nous vos retours sur l'application PyDesign.</p>
        <form id="contactForm" onsubmit="event.preventDefault(); document.getElementById('success').style.display='block';">
            <div class="input-group">
                <label>NOM COMPLET</label>
                <input type="text" placeholder="Alex Mercer" required>
            </div>
            <div class="input-group">
                <label>ADRESSE E-MAIL</label>
                <input type="email" placeholder="alex@example.com" required>
            </div>
            <div class="input-group">
                <label>MESSAGE</label>
                <textarea rows="4" placeholder="Votre message..." required></textarea>
            </div>
            <button type="submit">Envoyer le message</button>
            <div id="success" class="success-msg">✓ Message envoyé avec succès !</div>
        </form>
    </div>
</body>
</html>"""
    },
    "presentation": {
        "type": "python",
        "explanation": "Voici un script Python qui génère une présentation PowerPoint (.pptx) moderne et épurée (diapositive de titre, objectifs, et conclusion) en utilisant la bibliothèque python-pptx. Cliquez sur l'onglet **Preview** ou lancez le script pour générer et télécharger le fichier !",
        "code": """import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Initialisation de la présentation
prs = Presentation()
prs.slide_width = Inches(13.33)  # Format 16:9
prs.slide_height = Inches(7.5)

# Palette de couleurs "Sleek Dark"
COLOR_BG = RGBColor(11, 15, 25)       # Slate très sombre #0b0f19
COLOR_PRIMARY = RGBColor(99, 102, 241) # Indigo #6366f1
COLOR_TEXT = RGBColor(248, 250, 252)   # Blanc cassé #f8fafc
COLOR_MUTED = RGBColor(148, 163, 184)  # Slate secondaire #94a3b8

def set_slide_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

# --- Slide 1: Titre ---
slide_layout = prs.slide_layouts[6] # Mise en page vierge
slide1 = prs.slides.add_slide(slide_layout)
set_slide_background(slide1, COLOR_BG)

# Titre principal
title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.33), Inches(3.0))
tf = title_box.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "ATHENADESIGN STUDIO"
p.font.name = 'Arial'
p.font.size = Pt(54)
p.font.bold = True
p.font.color.rgb = COLOR_PRIMARY
p.space_after = Pt(14)

p2 = tf.add_paragraph()
p2.text = "Génération automatisée de présentations PowerPoint professionnelles"
p2.font.name = 'Arial'
p2.font.size = Pt(22)
p2.font.color.rgb = COLOR_TEXT

# --- Slide 2: Contenu / Objectifs ---
slide2 = prs.slides.add_slide(slide_layout)
set_slide_background(slide2, COLOR_BG)

# Titre de diapositive
title_box = slide2.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.33), Inches(1.0))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "Fonctionnalités Clés"
p.font.name = 'Arial'
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = COLOR_PRIMARY

# Zone de texte principale (points clés)
content_box = slide2.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.33), Inches(4.5))
tf_content = content_box.text_frame
tf_content.word_wrap = True

points = [
    ("Orchestration Multi-Agent", "Routage sémantique intelligent et collaboration organique entre agents spécialisés."),
    ("Visualisation interactive (AthenaDesign)", "Cockpit de télémétrie, bureau virtuel en 3D isométrique et annotations graphiques."),
    ("Génération PPTX Native", "Création automatique de livrables PowerPoint de qualité professionnelle directement depuis l'interface.")
]

for title, desc in points:
    p_title = tf_content.add_paragraph() if tf_content.text else tf_content.paragraphs[0]
    p_title.text = f"• {title}"
    p_title.font.name = 'Arial'
    p_title.font.size = Pt(22)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_TEXT
    p_title.space_after = Pt(4)
    
    p_desc = tf_content.add_paragraph()
    p_desc.text = f"  {desc}"
    p_desc.font.name = 'Arial'
    p_desc.font.size = Pt(16)
    p_desc.font.color.rgb = COLOR_MUTED
    p_desc.space_after = Pt(20)

# --- Slide 3: Conclusion ---
slide3 = prs.slides.add_slide(slide_layout)
set_slide_background(slide3, COLOR_BG)

title_box = slide3.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(11.33), Inches(3.0))
tf = title_box.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "Prêt à démarrer ?"
p.font.name = 'Arial'
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = COLOR_PRIMARY
p.space_after = Pt(10)

p2 = tf.add_paragraph()
p2.text = "Votre présentation est générée dans le dossier de sandbox de votre projet."
p2.font.name = 'Arial'
p2.font.size = Pt(18)
p2.font.color.rgb = COLOR_TEXT

# Sauvegarde de la présentation
output_filename = "presentation_athena.pptx"
prs.save(output_filename)
print(f"Présentation PPTX générée avec succès : {output_filename}")
"""
    }
}

_FILE_MARKER = re.compile(r"^[ \t]*===[ \t]*FILE:[ \t]*(.+?)[ \t]*===[ \t]*$", re.MULTILINE)


def _parse_multifile_blocks(text: str):
    """Projet MULTI-FICHIERS : segments précédés de `=== FILE: chemin ===`.
    Renvoie (prose_avant, [{path, content}]) ; ("", []) si aucun marqueur."""
    marks = list(_FILE_MARKER.finditer(text))
    if not marks:
        return "", []
    prose = text[:marks[0].start()].strip()
    files = []
    for i, m in enumerate(marks):
        path = m.group(1).strip().strip('"').strip("'").lstrip("/")
        end = marks[i + 1].start() if i + 1 < len(marks) else len(text)
        content = text[m.end():end]
        # Le DERNIER segment peut embarquer les blocs <suggestions>/<tweaks> finaux → on coupe.
        content = re.split(r"<(?:suggestions|tweaks)>", content, maxsplit=1, flags=re.IGNORECASE)[0]
        content = content.strip("\n")
        # Retire une fence ```lang … ``` éventuelle autour du contenu d'un fichier.
        content = re.sub(r"^[ \t]*```[a-zA-Z0-9]*[ \t]*\n", "", content)
        content = re.sub(r"\n?[ \t]*```[ \t]*$", "", content)
        if path and ".." not in path.split("/"):   # anti-traversée dès le parsing
            files.append({"path": path, "content": content})
    return prose, files


_TWEAK_LABELS = {
    "primary": "Couleur principale", "secondary": "Couleur secondaire",
    "accent": "Accent", "bg": "Fond", "background": "Fond", "surface": "Surface",
    "text": "Texte", "fg": "Texte", "foreground": "Texte", "border": "Bordure",
    "radius": "Arrondi", "border-radius": "Arrondi", "gap": "Espacement",
    "spacing": "Espacement", "font-size": "Taille du texte", "shadow": "Ombre",
}


def _derive_tweaks_from_code(code: str, files: list) -> list:
    """FILET déterministe : si le modèle n'a pas émis de bloc <tweaks>, on dérive les réglages
    dynamiques depuis les variables CSS `:root { --x: … }` du design lui-même (couleurs →
    color picker, longueurs px/rem/% → slider). Garantit un dock de personnalisation utile
    même quand le LLM l'oublie (fréquent en mode multi-fichiers). Max 6 contrôles."""
    css = code or ""
    for f in (files or []):
        css += "\n" + (f.get("content") or "")
    tweaks, seen = [], set()
    for block in re.findall(r":root\s*\{([^}]*)\}", css, re.IGNORECASE | re.DOTALL):
        for name, value in re.findall(r"(--[\w-]+)\s*:\s*([^;]+);", block):
            name, value = name.strip(), value.strip()
            if name in seen or len(tweaks) >= 6:
                continue
            key = name[2:].lower()
            label = _TWEAK_LABELS.get(key) or key.replace("-", " ").replace("_", " ").capitalize()
            if re.match(r"#[0-9a-fA-F]{3,8}$", value) or re.match(r"(rgb|rgba|hsl|hsla)\(", value, re.I):
                hexv = value if value.startswith("#") else "#6366f1"  # picker ne gère que le hex
                tweaks.append({"type": "color", "label": label, "name": name,
                               "values": "", "default": hexv})
                seen.add(name)
            else:
                m = re.match(r"(-?\d*\.?\d+)\s*(px|rem|em|%|vh|vw|pt)$", value)
                if m:
                    num, unit = float(m.group(1)), m.group(2)
                    lo = 0 if num >= 0 else int(num * 2)
                    hi = max(int(round(num * 2)), int(round(num)) + 8)
                    tweaks.append({"type": "slider", "label": label, "name": name,
                                   "values": f"{lo}{unit},{hi}{unit}",
                                   "default": (f"{num:g}")})
                    seen.add(name)
    return tweaks


def parse_artifact_response(text: str) -> dict:
    """Extrait {type, explanation, code} de la réponse du LLM, de façon ROBUSTE.

    Beaucoup de modèles (qwen3…) n'émettent PAS les balises <artifact_*> et répondent en
    « prose + bloc de code ». On gère ces cas pour ne JAMAIS coller la prose dans le code
    (sinon le `<!DOCTYPE html>` est précédé de texte → rendu cassé). Ordre des stratégies :
    balises explicites → bloc fencé ```lang → détection du début du vrai code → sinon code vide
    (la prose devient l'explication, pas du code)."""
    text = (text or "").strip()
    artifact_type = ""
    explanation = ""
    code = ""

    # 0) PROJET MULTI-FICHIERS (`=== FILE: chemin ===`) : on extrait tous les fichiers, l'ENTRÉE
    #    (index.html sinon 1er .html) sert d'aperçu/`code`, le reste est écrit sous design/.
    _prose_mf, _files = _parse_multifile_blocks(text)
    _entry_path = ""
    if _files:
        entry = (next((f for f in _files if f["path"].lower().endswith("index.html")), None)
                 or next((f for f in _files if f["path"].lower().endswith((".html", ".htm"))), None)
                 or _files[0])
        _entry_path = entry["path"]
        code = entry["content"]
        explanation = _prose_mf

    # 1)–4) Artefact MONO-BLOC — uniquement si AUCUN projet multi-fichiers n'a été extrait.
    if not _files:
        # 1) Balises explicites (si le modèle les respecte).
        t = re.search(r"<artifact_type>(.*?)</artifact_type>", text, re.DOTALL)
        if t:
            artifact_type = t.group(1).strip().lower()
        e = re.search(r"<artifact_explanation>(.*?)</artifact_explanation>", text, re.DOTALL)
        if e:
            explanation = e.group(1).strip()
        c = re.search(r"<artifact_code>(.*?)</artifact_code>", text, re.DOTALL)
        if c:
            code = c.group(1).strip()
        else:
            # 2) Bloc de code fencé ```lang … ``` (cas le plus fréquent). On prend le PLUS GROS
            #    bloc comme code ; la prose qui le précède devient l'explication.
            fences = list(re.finditer(r"```([a-zA-Z0-9]*)\s*\n(.*?)```", text, re.DOTALL))
            if fences:
                best = max(fences, key=lambda m: len(m.group(2)))
                code = best.group(2).strip()
                if not artifact_type and best.group(1):
                    artifact_type = best.group(1).strip().lower()
                if not explanation:
                    explanation = text[:best.start()].strip()
            else:
                # 3) Pas de fence : repérer le DÉBUT du vrai code (HTML ou Python) et couper.
                m = re.search(r"(<!DOCTYPE html|<html|<\?xml|<svg\b|^\s*import\s+\w|^\s*from\s+\w+\s+import|^\s*def\s+\w)",
                              text, re.IGNORECASE | re.MULTILINE)
                if m and m.start() > 0:
                    explanation = text[:m.start()].strip()
                    code = text[m.start():].strip()
                elif m:
                    code = text
                else:
                    # 4) Rien d'identifiable comme code → on ne dumpe PAS la prose dans le code.
                    explanation = text
                    code = ""

    # Nettoyage : retirer des fences résiduelles autour du code.
    if code.startswith("```"):
        code = re.sub(r"^```[a-zA-Z0-9]*\s*\n?", "", code)
        code = re.sub(r"\n?```$", "", code).strip()

    # Type : balise si fiable, sinon déduit du contenu.
    low = code.lower()
    _mermaid_kw = (r"^(flowchart|graph\s|sequencediagram|classdiagram|statediagram(-v2)?|"
                   r"erdiagram|gantt|mindmap|journey|pie\b|gitgraph|quadrantchart|timeline|requirementdiagram)")
    _is_mermaid = (artifact_type == "mermaid"
                   or bool(re.match(_mermaid_kw, code.strip().lower())))
    _is_react = (artifact_type in ("jsx", "tsx", "react")
                 or (not re.search(r"<!doctype|<html", low)
                     and re.search(r"\buse(state|effect|ref|memo|callback)\s*\(|react\.|from\s+['\"]react['\"]|export\s+default\s+function|=>\s*\(?\s*<[a-z]", code, re.I)))
    if artifact_type in ("python", "py"):
        atype = "python"
    elif _is_mermaid:
        atype = "mermaid"
    elif _is_react:
        atype = "react"
    elif artifact_type in ("html", "javascript", "js", "css", "xml", "svg"):
        atype = "html"
    elif re.search(r"<!doctype html|<html|<body|<div|<svg", low):
        atype = "html"
    elif re.search(r"\bimport\s+\w|\bdef\s+\w|from\s+pptx|matplotlib|plotly|\bprint\(", low):
        atype = "python"
    else:
        atype = "html"

    explanation = (explanation or "Voici votre design.").strip()
    # L'explication ne doit jamais embarquer tout le code (garde-fou d'affichage).
    if len(explanation) > 1200:
        explanation = explanation[:1200].rstrip() + "…"

    # Extraction des suggestions d'étapes suivantes depuis les balises XML
    suggestions = []
    sug_match = re.search(r"<suggestions>(.*?)</suggestions>", text, re.DOTALL | re.IGNORECASE)
    if sug_match:
        suggestions_raw = sug_match.group(1).strip().split("\n")
        for sug in suggestions_raw:
            sug = sug.strip().lstrip("-").lstrip("*").lstrip(" ").strip()
            if sug:
                suggestions.append(sug)

    # Extraction des tweaks (styles personnalisés dynamiques) depuis les balises XML
    tweaks = []
    tweaks_match = re.search(r"<tweaks>(.*?)</tweaks>", text, re.DOTALL | re.IGNORECASE)
    if tweaks_match:
        tweaks_raw = tweaks_match.group(1).strip().split("\n")
        for line in tweaks_raw:
            parts = [p.strip() for p in line.split("|")]
            if len(parts) >= 5:
                tweaks.append({
                    "type": parts[0],
                    "label": parts[1],
                    "name": parts[2],
                    "values": parts[3],
                    "default": parts[4]
                })

    # Filet : si le modèle n'a pas fourni de tweaks (fréquent en multi-fichiers), on les dérive
    # des variables CSS du design — seulement pour le web (html/react), pas python/mermaid.
    if not tweaks and atype in ("html", "react"):
        tweaks = _derive_tweaks_from_code(code, _files)

    return {
        "type": atype,
        "explanation": explanation,
        "code": code,
        "tweaks": tweaks,
        "suggestions": suggestions,
        "files": _files,          # projet multi-fichiers (vide = artefact mono-bloc)
        "entry": _entry_path,     # fichier d'aperçu du projet (ex. index.html)
    }

def react_scaffold(code: str) -> str:
    """Encapsule un composant React (JSX) dans une page autonome : React + ReactDOM + Babel
    standalone + Tailwind (CDN), puis rend <App/>. On retire import/export (React est global).
    Utilisé côté serveur (export PDF / partage / raw) ; le front a un équivalent pour l'aperçu live."""
    body = code or ""
    # Retrait des imports ES (y compris multi-lignes)
    body = re.sub(r"import\b[\s\S]*?from\s*['\"][^'\"]+['\"];?", "", body)
    body = re.sub(r"import\s*['\"][^'\"]+['\"];?", "", body)
    body = re.sub(r"(?m)^\s*export\s+default\s+", "", body)
    body = re.sub(r"(?m)^\s*export\s+", "", body)
    # Retrait de tout échafaudage/montage existant pour éviter les doubles déclarations
    body = re.sub(r"(?m)^\s*const\s+_C\s*=.*$", "", body)
    body = re.sub(r"(?m)^.*ReactDOM\s*\.\s*createRoot[^\n]*$", "", body)
    body = re.sub(r"(?m)^\s*ReactDOM\s*\.\s*render[^\n]*$", "", body)
    body = re.sub(r"(?m)^\s*\w+\s*\.\s*render\s*\(\s*<\s*App[^\n]*$", "", body)

    user_script = (
        "const {useState,useEffect,useRef,useMemo,useCallback,useReducer,useContext,Fragment}=React;\n"
        + body
        + "\nconst _C=(typeof App!==\"undefined\"?App:(typeof Component!==\"undefined\"?Component:"
        + "function(){return React.createElement(\"div\",{style:{padding:24}},\"Aucun composant App trouvé.\");}));\n"
        + "ReactDOM.createRoot(document.getElementById(\"root\")).render(React.createElement(_C));\n"
    )
    # Neutraliser un </script> littéral
    safe_user = user_script.replace("</script", "<\\/script").replace("</SCRIPT", "<\\/script")
    
    return (
        "<!DOCTYPE html><html lang=\"fr\"><head><meta charset=\"utf-8\">"
        "<meta name=\"viewport\" content=\"width=device-width, initial-scale=1.0\">"
        "<script crossorigin src=\"https://unpkg.com/react@18/umd/react.production.min.js\"></script>"
        "<script crossorigin src=\"https://unpkg.com/react-dom@18/umd/react-dom.production.min.js\"></script>"
        "<script src=\"https://unpkg.com/@babel/standalone/babel.min.js\"></script>"
        "<script src=\"https://cdn.tailwindcss.com\"></script>"
        "<style>body{margin:0;font-family:Inter,system-ui,sans-serif}"
        "#__err{display:none;padding:24px;color:#e11d48;background:#161616;white-space:pre-wrap;"
        "font:13px/1.6 ui-monospace,Menlo,Consolas,monospace}</style></head>"
        "<body><div id=\"root\"></div><pre id=\"__err\"></pre>"
        "<script type=\"text/plain\" id=\"__src\">" + safe_user + "</script>"
        "<script>(function(){"
        "function showErr(msg){var e=document.getElementById(\"__err\");if(e){e.style.display=\"block\";"
        "e.textContent=\"\\u26a0\\ufe0f Aperçu React — \"+msg;}"
        "try{window.parent.postMessage({type:\"iframe-log\",level:\"stderr\",message:msg},\"*\");}catch(_){}}"
        "window.addEventListener(\"DOMContentLoaded\",function(){"
        "if(typeof Babel===\"undefined\"){showErr(\"Babel non chargé (réseau ?).\");return;}"
        "var src=document.getElementById(\"__src\").textContent;"
        "var compiled;"
        "try{compiled=Babel.transform(src,{presets:[[\"react\",{runtime:\"classic\"}]]}).code;}"
        "catch(err){showErr(\"Erreur de compilation : \"+((err&&err.message)||err));return;}"
        "try{var s=document.createElement(\"script\");s.textContent=compiled;document.body.appendChild(s);}"
        "catch(err){showErr(\"Erreur d'exécution : \"+((err&&err.message)||err));}"
        "});})();</script></body></html>"
    )



def mermaid_scaffold(code: str) -> str:
    """Encapsule un diagramme Mermaid dans une page autonome (mermaid.js via CDN). Le code
    est ÉCHAPPÉ (<, >, &) car il vit dans <pre> ; mermaid lit le textContent (déséchappé)."""
    esc = (code or "").replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
    return (
        "<!DOCTYPE html><html lang=\"fr\"><head><meta charset=\"utf-8\">"
        "<meta name=\"viewport\" content=\"width=device-width, initial-scale=1.0\">"
        "<script src=\"https://cdn.jsdelivr.net/npm/mermaid@11/dist/mermaid.min.js\"></script>"
        "<style>body{margin:0;padding:24px;display:flex;justify-content:center;"
        "font-family:Inter,system-ui,sans-serif;background:#fff}.mermaid{max-width:100%}</style></head>"
        "<body><pre class=\"mermaid\">" + esc + "</pre>"
        "<script>try{mermaid.initialize({startOnLoad:true,theme:'default'});}catch(e){"
        "document.body.innerHTML='<p style=\\'color:#b91c1c\\'>Erreur Mermaid: '+e.message+'</p>';}</script>"
        "</body></html>"
    )


def _athena_default_model() -> str:
    """Modèle d'AthenaDesign = LE MÊME CHOIX que le reste d'Athena (pas de knob dédié).
    On part du modèle de l'orchestrateur ; `swarm._complete` applique ensuite l'override
    `LLM_MODEL` de la config utilisateur → le choix de modèle global s'applique tel quel."""
    try:
        from core.state import swarm as _sw
        orch = getattr(_sw, "orchestrator_name", "Athena")
        agent = _sw.agents.get(orch)
        if agent and getattr(agent, "model", None):
            return agent.model
    except Exception:
        pass
    import os
    return os.getenv("DEFAULT_MODEL", "").strip() or "qwen3"


def _model_supports_vision(model: str) -> bool:
    """Vrai si le modèle accepte des images (multimodal). S'appuie sur litellm."""
    try:
        import litellm
        return bool(litellm.supports_vision(model=model))
    except Exception:
        return False


def _describe_images(images: list) -> str:
    """Repli SANS vision côté modèle principal : si un VISION_MODEL est configuré, on lui
    fait décrire les images → texte injecté comme contexte. Sinon chaîne vide."""
    import os
    vmodel = os.getenv("VISION_MODEL", "").strip()
    if not vmodel or not images:
        return ""
    from core.state import swarm as _sw
    out = []
    for i, data_url in enumerate(images[:4], 1):
        try:
            msg = [{"role": "user", "content": [
                {"type": "text", "text": "Décris précisément cette image de référence design "
                 "(couleurs, typographie, mise en page, composants, ambiance) pour qu'un autre "
                 "modèle puisse s'en inspirer. Retranscris tout texte visible."},
                {"type": "image_url", "image_url": {"url": data_url}},
            ]}]
            resp = _sw._complete(vmodel, msg, tools_schema=None, allow_continuation=False, allow_fallback=False)
            desc = (resp.choices[0].message.content or "").strip()
            if desc:
                out.append(f"[Image de référence {i}]\n{desc}")
        except Exception:
            pass
    return "\n\n".join(out)


def _build_system(design_system: str = "", context_text: str = "", note: str = "",
                  base_code: str = "") -> str:
    """Assemble le prompt système : règles AthenaDesign + charte (design system) + CODE
    EXISTANT du projet (base de travail) + contexte importé + note éventuelle."""
    parts = [SYSTEM_PROMPT]
    if (design_system or "").strip():
        parts.append("=== DESIGN SYSTEM (charte à RESPECTER impérativement : couleurs, "
                     "typographie, composants, ton) ===\n" + design_system.strip())
    if (base_code or "").strip():
        parts.append(
            "=== CODE ACTUEL DU PROJET (point de départ OBLIGATOIRE) ===\n"
            "Voici le code EXISTANT de la page du projet. Sauf demande explicite de repartir "
            "de zéro, tu dois PARTIR DE CE CODE : conserve sa structure, son contenu, son "
            "design ET ses fonctionnalités, et applique UNIQUEMENT la modification demandée "
            "DESSUS (n'invente PAS une nouvelle page, ne change PAS le style existant sauf si "
            "c'est explicitement demandé). Si le code est un projet MULTI-FICHIERS "
            "(`=== FILE: chemin ===`), réémets TOUS les fichiers dans CE MÊME format en "
            "préservant l'arborescence ; sinon produis la page COMPLÈTE modifiée.\n\n"
            + base_code.strip())
    if (context_text or "").strip():
        parts.append("=== CONTEXTE FOURNI PAR L'UTILISATEUR (références, documents, capture "
                     "web — inspire-t'en, ne recopie pas aveuglément) ===\n" + context_text.strip())
    if (note or "").strip():
        parts.append(note.strip())
    return "\n\n".join(parts)


def _generate_via_athena(prompt: str, history: list, model_name: str = "",
                         design_system: str = "", context_text: str = "", images: list = None,
                         base_code: str = "", on_delta=None) -> dict:
    """Génère via l'INFRA LLM d'Athena (swarm._complete) : endpoint/clés/fallback d'Athena.
    Gère la charte (design_system), le contexte importé (docs/web) et les IMAGES :
    modèle vision direct → sinon pré-description via VISION_MODEL → sinon note (marche au max
    sans vision). Synchrone (à appeler en thread)."""
    from core.state import swarm as _sw
    images = images or []
    # Modèle Design : choix UI explicite > DESIGN_MODEL (réglage par-user) > modèle du chat /
    # défaut orchestrateur. `_explicit` = un modèle dédié a été choisi → on le FORCE (sinon
    # l'override global LLM_MODEL dans _complete l'écraserait) ; vide = la feature suit le chat.
    _design_pref = ""
    try:
        from core import user_config
        _design_pref = (user_config.get_all().get("DESIGN_MODEL") or "").strip()
    except Exception:
        _design_pref = ""
    _explicit = (model_name or "").strip() or _design_pref
    model = _explicit or _athena_default_model()

    note = ""
    user_content = prompt
    if images:
        if _model_supports_vision(model):
            # Modèle multimodal : on envoie les images directement.
            user_content = [{"type": "text", "text": prompt}] + [
                {"type": "image_url", "image_url": {"url": u}} for u in images[:6]]
        else:
            desc = _describe_images(images)
            if desc:
                context_text = (context_text + "\n\n" + desc).strip()
            else:
                note = (f"[{len(images)} image(s) de référence jointe(s) mais le modèle n'est pas "
                        "multimodal et aucun VISION_MODEL n'est configuré → non analysées. "
                        "Demande à l'utilisateur de les décrire, ou configure un modèle vision.]")

    messages = [{"role": "system", "content": _build_system(design_system, context_text, note, base_code)}]
    for m in (history or []):
        role = m.get("role")
        if role in ("user", "assistant") and m.get("content"):
            messages.append({"role": role, "content": m["content"]})
    messages.append({"role": "user", "content": user_content})
    # Budget de sortie ÉLEVÉ : un design (HTML/React complet) dépasse vite 4000 tokens → sinon
    # le code est TRONQUÉ (balise/JSX coupés → aperçu cassé). + auto-continuation en filet.
    _mt = int(os.getenv("ATHENADESIGN_MAX_TOKENS", "8192") or 8192)
    # FORCE le modèle Design dédié pour ce run (prime sur LLM_MODEL) ; sinon laisse le chat décider.
    from core.state import _forced_model
    _tok = _forced_model.set(model) if _explicit else None
    try:
        resp = _sw._complete(model, messages, tools_schema=None, allow_continuation=True,
                             allow_fallback=True, max_tokens=_mt, on_delta=on_delta)
    finally:
        if _tok is not None:
            _forced_model.reset(_tok)
    text = (resp.choices[0].message.content or "")
    _u = getattr(resp, "usage", None)
    return _attach_usage(parse_artifact_response(text),
                         getattr(_u, "prompt_tokens", 0), getattr(_u, "completion_tokens", 0))


def _attach_usage(result: dict, prompt_tokens, completion_tokens) -> dict:
    """Ajoute la consommation de tokens au résultat (clé `usage`) — affiché côté studio
    et remonté à la télémétrie globale. Tolère les valeurs manquantes (0)."""
    pt = int(prompt_tokens or 0)
    ct = int(completion_tokens or 0)
    result["usage"] = {"prompt_tokens": pt, "completion_tokens": ct, "total_tokens": pt + ct}
    return result


async def generate_design(prompt: str, history: list, provider: str = "athena",
                          api_key: str = "", model_name: str = "",
                          design_system: str = "", context_text: str = "", images: list = None,
                          base_code: str = "") -> dict:
    """Routeur LLM. Par DÉFAUT (provider 'athena' ou non précisé), passe par l'infra LLM
    d'Athena (clés/endpoint/fallback configurés). 'mock' → templates hors-ligne. Les providers
    externes explicites (gemini/anthropic/openai) restent possibles si une clé est fournie.
    design_system/context_text/images : charte + contexte importé + images de référence."""
    import asyncio

    if provider == "mock":
        prompt_lower = prompt.lower()
        if "simulation" in prompt_lower or "chaos" in prompt_lower or "lorenz" in prompt_lower:
            return MOCK_TEMPLATES["simulation"]
        elif "plot" in prompt_lower or "graph" in prompt_lower or "sinus" in prompt_lower or "pandas" in prompt_lower:
            return MOCK_TEMPLATES["plot"]
        elif "form" in prompt_lower or "contact" in prompt_lower or "feedback" in prompt_lower:
            return MOCK_TEMPLATES["form"]
        elif "powerpoint" in prompt_lower or "presentation" in prompt_lower or "pptx" in prompt_lower or "slide" in prompt_lower:
            return MOCK_TEMPLATES["presentation"]
        else:
            # Default to dashboard template
            return MOCK_TEMPLATES["dashboard"]

    # Par DÉFAUT : infra LLM d'Athena (swarm._complete) — endpoint/clés/fallback configurés
    # dans Athena. On ne tombe sur un provider EXTERNE que s'il est explicitement demandé
    # ET qu'une clé est fournie par le front.
    if provider not in ("gemini", "anthropic", "openai") or not api_key:
        try:
            return await asyncio.to_thread(
                _generate_via_athena, prompt, history, model_name, design_system, context_text, images, base_code)
        except Exception as e:
            return {"type": "html",
                    "explanation": f"⚠️ Erreur de génération via l'API d'Athena : {e}",
                    "code": ""}

    # Real LLM API calls (providers externes explicites avec clé fournie)
    headers = {}
    payload = {}
    url = ""
    
    # Budget de sortie élevé (design long) + timeout généreux (génération longue) : sinon
    # le code est TRONQUÉ (limite tokens) ou coupé par timeout → aperçu cassé/blanc.
    _mt = int(os.getenv("ATHENADESIGN_MAX_TOKENS", "8192") or 8192)
    try:
        async with httpx.AsyncClient(timeout=180.0) as client:

            if provider == "gemini":
                # Using Gemini API (Google AI Studio model, e.g. gemini-2.5-flash)
                # Google standard endpoint: /v1beta/models/{model}:generateContent
                model = model_name or "gemini-2.5-flash"
                url = f"https://generativelanguage.googleapis.com/v1beta/models/{model}:generateContent?key={api_key}"
                
                # Format conversation history
                contents = []
                # System prompt as systemInstruction if supported, or prepended to the first message
                # For simplicity, we can pass systemInstruction inside systemInstruction parameter
                
                for msg in history:
                    contents.append({
                        "role": "user" if msg["role"] == "user" else "model",
                        "parts": [{"text": msg["content"]}]
                    })
                
                # Add current user prompt
                contents.append({
                    "role": "user",
                    "parts": [{"text": prompt}]
                })
                
                payload = {
                    "contents": contents,
                    "systemInstruction": {
                        "parts": [{"text": SYSTEM_PROMPT}]
                    },
                    "generationConfig": {
                        "temperature": 0.2,
                        "maxOutputTokens": _mt
                    }
                }
                
                resp = await client.post(url, json=payload)
                resp.raise_for_status()
                data = resp.json()
                
                text_out = data["candidates"][0]["content"]["parts"][0]["text"]
                _um = data.get("usageMetadata", {}) or {}
                return _attach_usage(parse_artifact_response(text_out),
                                     _um.get("promptTokenCount", 0), _um.get("candidatesTokenCount", 0))
                
            elif provider == "anthropic":
                # Using Anthropic API
                model = model_name or "claude-3-5-sonnet-latest"
                url = "https://api.anthropic.com/v1/messages"
                headers = {
                    "x-api-key": api_key,
                    "anthropic-version": "2023-06-01",
                    "content-type": "application/json"
                }
                
                # Format conversation history
                messages = []
                for msg in history:
                    if msg["role"] in ["user", "assistant"]:
                        messages.append({
                            "role": msg["role"],
                            "content": msg["content"]
                        })
                messages.append({
                    "role": "user",
                    "content": prompt
                })
                
                payload = {
                    "model": model,
                    "max_tokens": _mt,
                    "system": SYSTEM_PROMPT,
                    "messages": messages,
                    "temperature": 0.2
                }
                
                resp = await client.post(url, json=payload, headers=headers)
                resp.raise_for_status()
                data = resp.json()
                
                text_out = data["content"][0]["text"]
                _u = data.get("usage", {}) or {}
                return _attach_usage(parse_artifact_response(text_out),
                                     _u.get("input_tokens", 0), _u.get("output_tokens", 0))
                
            elif provider == "openai":
                # Using OpenAI API
                model = model_name or "gpt-4o-mini"
                url = "https://api.openai.com/v1/chat/completions"
                headers = {
                    "Authorization": f"Bearer {api_key}",
                    "Content-Type": "application/json"
                }
                
                messages = [{"role": "system", "content": SYSTEM_PROMPT}]
                for msg in history:
                    messages.append({
                        "role": msg["role"],
                        "content": msg["content"]
                    })
                messages.append({
                    "role": "user",
                    "content": prompt
                })
                
                payload = {
                    "model": model,
                    "messages": messages,
                    "temperature": 0.2,
                    "max_tokens": _mt
                }

                resp = await client.post(url, json=payload, headers=headers)
                resp.raise_for_status()
                data = resp.json()

                text_out = data["choices"][0]["message"]["content"]
                _u = data.get("usage", {}) or {}
                return _attach_usage(parse_artifact_response(text_out),
                                     _u.get("prompt_tokens", 0), _u.get("completion_tokens", 0))
                
            else:
                raise ValueError(f"Provider inconnu : {provider}")
                
    except Exception as e:
        # Fallback to Mock template on network/API failure so the user isn't stuck
        fallback = MOCK_TEMPLATES["dashboard"]
        fallback["explanation"] = f"⚠️ [Erreur API : {str(e)}]\nUne erreur est survenue lors de l'appel à {provider.capitalize()}. Chargement du modèle de démonstration local à la place."
        return fallback
