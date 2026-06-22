"""AthenaDesign — exécution du code généré, ISOLÉE par défaut.

Le code de présentation/visualisation produit par le LLM est exécuté :
  - EN PRIORITÉ dans un conteneur Docker (réseau coupé, --cap-drop ALL, racine en lecture
    seule hors /work) via tools.sandbox_runner.run_python_in_dir, avec une image dérivée
    contenant python-pptx / matplotlib / numpy / pandas / plotly ;
  - EN REPLI (Docker indisponible OU SANDBOX_MODE=off) : subprocess local — non isolé,
    risque assumé et JOURNALISÉ (identique à l'ancien comportement).

Variables d'environnement :
  ATHENADESIGN_DOCKER_IMAGE  image Docker à utiliser telle quelle (sinon image dérivée
                             auto, construite une fois et cachée par tag)
  ATHENADESIGN_PIP           libs pip de l'image dérivée (défaut ci-dessous)
  ATHENADESIGN_TIMEOUT       budget d'exécution en secondes (défaut 30)
  SANDBOX_MODE=off           force le repli local non isolé
"""
import os
import sys
import subprocess
import time
import shutil
import re
import threading
import logging

logger = logging.getLogger("athenadesign")

SANDBOX_DIR = os.path.abspath(os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "sandbox"))

_DESIGN_TAG = "athena-design:latest"
_DESIGN_PIP_DEFAULT = "python-pptx matplotlib numpy pandas plotly"
_build_lock = threading.Lock()
_build_failed = False  # mémorise un échec de build pour ne pas le retenter en boucle


def prepare_sandbox(project_id: str) -> str:
    """Crée un dossier sandbox propre pour l'exécution du projet."""
    project_sandbox = os.path.join(SANDBOX_DIR, project_id)
    if os.path.exists(project_sandbox):
        shutil.rmtree(project_sandbox)
    os.makedirs(project_sandbox, exist_ok=True)
    return project_sandbox


# Patches injectés AVANT le code généré (capture matplotlib/plotly). Chaîne normale (pas
# f-string) → pas d'échappement d'accolades.
_PRE_PATCHES = """# --- AUTOMATIC ATHENADESIGN SANDBOX PATCHES ---
import sys, os, io

try:
    import matplotlib
    matplotlib.use('Agg')
    import matplotlib.pyplot as plt
    def _patched_show(*args, **kwargs):
        fig_num = len([f for f in os.listdir('.') if f.startswith('plot_') and f.endswith('.png')]) + 1
        filename = f"plot_{fig_num}.png"
        plt.savefig(filename, bbox_inches='tight', dpi=150)
        plt.close()
        print(f"[AthenaDesign Plot Saved: {filename}]")
    plt.show = _patched_show
except ImportError:
    pass

try:
    import plotly.io as pio
    from plotly.io._renderers import ExternalRenderer
    class AthenaDesignRenderer(ExternalRenderer):
        def render(self, fig):
            fig_num = len([f for f in os.listdir('.') if f.startswith('plotly_') and f.endswith('.html')]) + 1
            filename = f"plotly_{fig_num}.html"
            fig.write_html(filename)
            print(f"[AthenaDesign Interactive Plot Saved: {filename}]")
    pio.renderers['athenadesign'] = AthenaDesignRenderer()
    pio.renderers.default = 'athenadesign'
except Exception:
    pass
# --- END OF PATCHES ---
"""

# Patches APRÈS le code : fallback matplotlib + ANTI-DÉBORDEMENT PPTX déterministe.
_POST_PATCHES = """
# Figures matplotlib non sauvegardées : fallback
try:
    if 'plt' in locals() or 'plt' in globals():
        import matplotlib.pyplot as plt
        if plt.get_fignums():
            plt.show()
except Exception:
    pass

# Anti-débordement PowerPoint (déterministe, indépendant du modèle) : pour tout .pptx
# produit, on force le texte à RENTRER dans sa boîte (word_wrap + shrink-to-fit) et on
# ramène les formes qui sortent dans les limites de la diapo.
try:
    import glob as _glob
    from pptx import Presentation as _Prs
    from pptx.enum.text import MSO_AUTO_SIZE as _AS
    for _f in _glob.glob('*.pptx'):
        try:
            _p = _Prs(_f)
            _sw, _sh = _p.slide_width, _p.slide_height
            for _sl in _p.slides:
                _to_remove = []
                for _shp in _sl.shapes:
                    # Formes ABERRANTES (dimension <= 0 : ex. « barres » dessinées en hauteur
                    # négative) → on les SUPPRIME (sinon rendu cassé/illisible).
                    try:
                        if (_shp.width is not None and _shp.width <= 0) or \
                           (_shp.height is not None and _shp.height <= 0):
                            _to_remove.append(_shp)
                            continue
                    except Exception:
                        pass
                    if _shp.has_text_frame:
                        _tf = _shp.text_frame
                        _tf.word_wrap = True
                        try:
                            _tf.auto_size = _AS.TEXT_TO_FIT_SHAPE
                        except Exception:
                            pass
                    try:
                        # Clamp des formes plus GRANDES que la diapo, puis repositionnement.
                        if _shp.width and _shp.width > _sw:
                            _shp.width = _sw
                        if _shp.height and _shp.height > _sh:
                            _shp.height = _sh
                        if _shp.left is not None and _shp.width and _shp.left + _shp.width > _sw:
                            _shp.left = max(0, _sw - _shp.width)
                        if _shp.top is not None and _shp.height and _shp.top + _shp.height > _sh:
                            _shp.top = max(0, _sh - _shp.height)
                        if _shp.left is not None and _shp.left < 0:
                            _shp.left = 0
                        if _shp.top is not None and _shp.top < 0:
                            _shp.top = 0
                    except Exception:
                        pass
                for _shp in _to_remove:
                    try:
                        _shp._element.getparent().remove(_shp._element)
                    except Exception:
                        pass
                # LISIBILITÉ : on remet TOUT shape porteur de TEXTE au PREMIER PLAN (les formes
                # décoratives colorées, ajoutées après, passaient DEVANT le texte → illisible).
                # Réordonner le spTree = appendre les éléments texte à la fin (= dessus).
                try:
                    _spTree = _sl.shapes._spTree
                    for _shp in list(_sl.shapes):
                        if _shp.has_text_frame and (_shp.text_frame.text or "").strip():
                            _spTree.append(_shp._element)
                except Exception:
                    pass
            _p.save(_f)
            print(f"[AthenaDesign PPTX ajuste: {_f}]")
        except Exception as _e:
            print(f"[AthenaDesign PPTX post-traitement ignore: {_e}]")
except Exception:
    pass
"""


def _patched_code(code: str) -> str:
    """Encadre le code généré : patches de capture (avant) + fallback & anti-débordement
    pptx (après). Les fichiers produits (plots, .pptx ajusté) restent dans le sandbox."""
    return _PRE_PATCHES + "\n" + code + "\n" + _POST_PATCHES


def _image_exists(tag: str) -> bool:
    try:
        r = subprocess.run(["docker", "image", "inspect", tag], capture_output=True, timeout=15)
        return r.returncode == 0
    except Exception:
        return False


def _ensure_design_image():
    """Renvoie (tag, erreur). Si ATHENADESIGN_DOCKER_IMAGE est défini, on l'utilise tel quel.
    Sinon, on construit UNE FOIS une image dérivée (base sandbox + libs design), cachée."""
    global _build_failed
    explicit = os.getenv("ATHENADESIGN_DOCKER_IMAGE", "").strip()
    if explicit:
        return explicit, None
    if _image_exists(_DESIGN_TAG):
        return _DESIGN_TAG, None
    if _build_failed:
        return None, "build précédent échoué"
    with _build_lock:
        if _image_exists(_DESIGN_TAG):
            return _DESIGN_TAG, None
        base = os.getenv("SANDBOX_DOCKER_IMAGE", "python:3.13-slim")
        pip_pkgs = os.getenv("ATHENADESIGN_PIP", _DESIGN_PIP_DEFAULT)
        dockerfile = f"FROM {base}\nRUN pip install --no-cache-dir {pip_pkgs}\n"
        logger.info("AthenaDesign : construction de l'image sandbox %s (libs: %s)…", _DESIGN_TAG, pip_pkgs)
        try:
            r = subprocess.run(
                ["docker", "build", "-q", "-t", _DESIGN_TAG, "-"],
                input=dockerfile, capture_output=True, text=True, timeout=600,
            )
            if r.returncode != 0:
                _build_failed = True
                return None, (r.stderr or r.stdout or "échec docker build")[-500:]
            return _DESIGN_TAG, None
        except Exception as e:
            _build_failed = True
            return None, str(e)


def _scan_outputs(sandbox_path: str):
    """Recense les fichiers produits dans le sandbox (plots, html interactifs, autres)."""
    plots, interactive_plots, other_files = [], [], []
    if os.path.exists(sandbox_path):
        for filename in os.listdir(sandbox_path):
            if filename == "run.py":
                continue
            filepath = os.path.join(sandbox_path, filename)
            if os.path.isfile(filepath):
                if filename.startswith("plot_") and filename.endswith(".png"):
                    plots.append(filename)
                elif filename.startswith("plotly_") and filename.endswith(".html"):
                    interactive_plots.append(filename)
                else:
                    other_files.append({"name": filename, "size": os.path.getsize(filepath)})
    _key = lambda x: [int(s) if s.isdigit() else s for s in re.split(r'(\d+)', x)]
    plots.sort(key=_key)
    interactive_plots.sort(key=_key)
    return plots, interactive_plots, other_files


def _run_local(patched_code: str, sandbox_path: str, timeout: int):
    """REPLI non isolé : subprocess local (droits du serveur). Utilisé seulement si
    Docker indisponible ou SANDBOX_MODE=off. Journalisé comme un risque."""
    script_path = os.path.join(sandbox_path, "run.py")
    with open(script_path, "w", encoding="utf-8") as f:
        f.write(patched_code)
    try:
        process = subprocess.Popen(
            [sys.executable, "run.py"], cwd=sandbox_path,
            stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True, encoding="utf-8",
        )
        stdout, stderr = process.communicate(timeout=timeout)
        return stdout, stderr, (process.returncode == 0)
    except subprocess.TimeoutExpired as e:
        process.kill()
        return (e.stdout or ""), (e.stderr or "") + f"\n[Execution Error: Timeout ({timeout}s)]", False
    except Exception as e:
        return "", f"[Execution Error: {e}]", False


def execute_code(code: str, project_id: str) -> dict:
    """Exécute le code généré et renvoie stdout/stderr + fichiers produits. Isolé via Docker
    par défaut, repli local non isolé sinon (journalisé). `sandboxed` indique le mode réel."""
    sandbox_path = prepare_sandbox(project_id)
    patched = _patched_code(code)
    timeout = int(os.getenv("ATHENADESIGN_TIMEOUT", "30") or 30)
    start_time = time.time()
    sandboxed = False

    try:
        from tools import sandbox_runner
        _docker_ok = sandbox_runner.sandbox_mode() != "off" and sandbox_runner.docker_available()
    except Exception:
        _docker_ok = False

    if _docker_ok:
        image, build_err = _ensure_design_image()
        if image:
            stdout, stderr, rc = sandbox_runner.run_python_in_dir(
                patched, sandbox_path, image=image, timeout=timeout)
            success = (rc == 0)
            sandboxed = True
        else:
            logger.warning("AthenaDesign : image sandbox indisponible (%s) → exécution LOCALE non isolée.", build_err)
            stdout, stderr, success = _run_local(patched, sandbox_path, timeout)
    else:
        if os.getenv("SANDBOX_MODE", "docker").strip().lower() != "off":
            logger.warning("AthenaDesign : Docker indisponible → exécution LOCALE non isolée (risque).")
        stdout, stderr, success = _run_local(patched, sandbox_path, timeout)

    elapsed_time = time.time() - start_time
    plots, interactive_plots, other_files = _scan_outputs(sandbox_path)
    return {
        "success": success,
        "stdout": stdout,
        "stderr": stderr,
        "plots": plots,
        "interactive_plots": interactive_plots,
        "other_files": other_files,
        "execution_time": round(elapsed_time, 3),
        "sandboxed": sandboxed,
    }
